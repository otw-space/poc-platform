import os
import json
import re
import uuid as uuid_lib
from datetime import datetime
from typing import Any
from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from ..database import get_db
from ..models.sop import SopDocument, TestCaseCategory, TestCase, ScriptFile
from ..models.user import User
from ..schemas.sop import (
    SopDocumentCreate, SopDocumentUpdate, SopDocumentOut,
    TestCaseCategoryCreate, TestCaseCategoryUpdate, TestCaseCategoryOut,
    TestCaseCreate, TestCaseUpdate, TestCaseOut, TestCaseListOut,
    ScriptCreate, ScriptUpdate, ScriptOut,
)
from ..middleware.auth import get_current_user, require_permission
from ..services.logger import log_operation

router = APIRouter(prefix="/api/sops", tags=["sops"])

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "uploads", "sops")
MAX_FILE_SIZE = 50 * 1024 * 1024

ALLOWED_DOC_EXT = {".doc", ".docx", ".pdf"}
ALLOWED_REPORT_EXT = {".ppt", ".pptx", ".pdf"}

def _resolve_image_path(src: str) -> str | None:
    """Resolve an image URL to a local file path. Supports /api/sops/documents/{id}/images/{name} URLs."""
    if not src:
        return None
    # Relative or absolute server URL
    import re
    m = re.match(r'(?:/api)?/sops/documents/([^/]+)/images/([^/\s]+)', src)
    if m:
        doc_id, filename = m.group(1), m.group(2)
        path = os.path.join(UPLOAD_DIR, "documents", doc_id, "images", filename)
        if os.path.exists(path):
            return path
    # Absolute file path
    if os.path.exists(src):
        return src
    # Try resolving relative to upload dir
    return None


CJK_FONT_PATH = "/System/Library/Fonts/STHeiti Light.ttc"
if not os.path.exists(CJK_FONT_PATH):
    CJK_FONT_PATH = "/System/Library/Fonts/Supplemental/Songti.ttc"
if not os.path.exists(CJK_FONT_PATH):
    CJK_FONT_PATH = None


def _ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)
    return path


def _parse_table_rows(content: str) -> list[list[str]] | None:
    """Parse markdown tables (with tabs, <br>, spacing variations). Returns rows or None."""
    lines = content.strip().split('\n')
    rows = []
    found_separator = False
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith('|'):
            return None  # Not a table
        # Check for separator line like |:---:|:---|
        if re.match(r'^[\|\s\-:]+$', stripped):
            found_separator = True
            continue
        # Parse cells: split by |, handle tab-separated cells
        cells = [c.strip() for c in stripped.split('|')[1:-1]]
        if cells:
            rows.append(cells)
    return rows if rows else None  # Allow tables without separator lines


def _walk_html(el, tokens: list):
    """Recursively walk HTML elements, appending structured tokens in document order."""
    BLOCK_TAGS = {'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol',
                  'table', 'pre', 'code', 'blockquote', 'div', 'body'}
    if not hasattr(el, 'name'):
        return
    name = (el.name or '').lower()

    # Handle <img> tags — produce image token
    if name == 'img':
        src = el.get('src', '') if hasattr(el, 'get') else ''
        alt = el.get('alt', '') if hasattr(el, 'get') else ''
        if src:
            tokens.append({'type': 'image', 'text': alt or 'image', 'level': 0, 'items': [], 'src': src})
        return

    if name not in BLOCK_TAGS:
        if hasattr(el, 'children'):
            for child in el.children:
                _walk_html(child, tokens)
        return

    # For block tags, check if they contain images before checking text
    imgs = el.find_all('img') if hasattr(el, 'find_all') else []
    text = el.get_text(strip=True) if hasattr(el, 'get_text') else ''

    # Emit images found inside this block
    for img in imgs:
        src = img.get('src', '') if hasattr(img, 'get') else ''
        alt = img.get('alt', '') if hasattr(img, 'get') else ''
        if src:
            tokens.append({'type': 'image', 'text': alt or 'image', 'level': 0, 'items': [], 'src': src})

    # Remove image alt text from block text (markdown puts alt text in p)
    if imgs:
        for img in imgs:
            alt = img.get('alt', '') if hasattr(img, 'get') else ''
            if alt:
                text = text.replace(alt, '')
        text = text.strip()

    if not text and name not in ('ul', 'ol', 'table'):
        return

    token = {'type': name, 'text': text, 'level': 1, 'items': []}
    if name in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
        token['level'] = int(name[1])
    elif name in ('ul', 'ol'):
        token['items'] = [li.get_text(strip=True) for li in el.find_all('li', recursive=False)]
    elif name in ('pre', 'code'):
        token['type'] = 'code'
        token['text'] = el.get_text()
    elif name == 'blockquote':
        pass
    elif name == 'table':
        rows = []
        for row in el.find_all('tr'):
            rows.append([c.get_text(strip=True) for c in row.find_all(['td', 'th'])])
        token['type'] = 'table'
        token['items'] = rows
    elif name in ('div', 'body'):
        if hasattr(el, 'children'):
            for child in el.children:
                _walk_html(child, tokens)
        return
    tokens.append(token)


def _normalize_markdown(content: str) -> str:
    """Pre-process markdown to fix tables missing separator lines.
    Only fixes the first row of a table (header) that lacks a following separator.
    Does NOT modify rows that already have valid separators."""
    lines = content.split('\n')
    result = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        # Detect the start of a potential table (first pipe-delimited row)
        is_pipe_row = (stripped.startswith('|') and stripped.endswith('|')
                       and '---' not in stripped.replace(' ', '')
                       and not re.match(r'^[\|\s\-:]+$', stripped))
        if is_pipe_row:
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            if len(cells) >= 2:
                # Look at the previous line — was it already part of a table?
                # If this pipe row follows a pipe row or separator, it's a continuation, don't touch
                prev_non_empty = ''
                for r in reversed(result):
                    if r.strip():
                        prev_non_empty = r.strip()
                        break
                is_continuation = prev_non_empty.startswith('|')
                if is_continuation:
                    result.append(lines[i])
                    i += 1
                    continue
                # This is the first row of a table block — ensure blank line before
                if result and result[-1].strip() and not result[-1].strip().startswith('|'):
                    result.append('')
                result.append(lines[i])
                # Check if next non-empty line is a separator
                next_idx = i + 1
                while next_idx < len(lines) and not lines[next_idx].strip():
                    next_idx += 1
                if next_idx < len(lines) and not re.match(r'^[\|\s\-:]+$', lines[next_idx].strip()):
                    # Next line is NOT a separator — insert one
                    sep = '|' + '|'.join(['---' for _ in range(len(cells))]) + '|'
                    result.append(sep)
                i += 1
                continue
        result.append(lines[i])
        i += 1
    return '\n'.join(result)


def _render_markdown(content: str) -> list[dict]:
    """Render markdown to structured tokens.
    Returns list of {type, text, level, items}.
    Uses markdown + BeautifulSoup if available, with recursive tree walking.
    """
    try:
        import markdown as md_lib
        from bs4 import BeautifulSoup
        content = _normalize_markdown(content)
        html = md_lib.markdown(content, extensions=['extra', 'codehilite', 'tables'])
        soup = BeautifulSoup(html, 'html.parser')
        tokens = []
        if hasattr(soup, 'children'):
            for child in soup.children:
                _walk_html(child, tokens)
        if tokens:
            return tokens
    except ImportError:
        pass

    # Fallback parser (also used as primary if markdown lib not available)
    tokens = []
    in_code = False
    code_lines = []
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Code block
        if stripped.startswith('```'):
            if in_code:
                tokens.append({'type': 'code', 'text': '\n'.join(code_lines), 'level': 0, 'items': []})
                code_lines = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(line)
            i += 1
            continue

        # Try table parsing
        if stripped.startswith('|') and stripped.endswith('|'):
            # Collect contiguous table lines
            table_block = []
            j = i
            while j < len(lines):
                lj = lines[j].strip()
                if lj.startswith('|'):
                    table_block.append(lines[j])
                    j += 1
                elif not lj:  # empty line doesn't break if next line is table
                    if j + 1 < len(lines) and lines[j + 1].strip().startswith('|'):
                        j += 1
                        continue
                    else:
                        break
                else:
                    break
            rows = _parse_table_rows('\n'.join(table_block))
            if rows:
                tokens.append({'type': 'table', 'text': '', 'level': 0, 'items': rows})
                i = j
                continue
            # Not a valid table, fall through to paragraph
            tokens.append({'type': 'p', 'text': stripped, 'level': 0, 'items': []})
            i += 1
            continue

        # Headings
        if stripped.startswith('# ') and len(stripped) > 2:
            tokens.append({'type': 'h1', 'text': stripped[2:], 'level': 1, 'items': []})
        elif stripped.startswith('## ') and len(stripped) > 3:
            tokens.append({'type': 'h2', 'text': stripped[3:], 'level': 2, 'items': []})
        elif stripped.startswith('### ') and len(stripped) > 4:
            tokens.append({'type': 'h3', 'text': stripped[4:], 'level': 3, 'items': []})
        elif stripped.startswith('#### ') and len(stripped) > 5:
            tokens.append({'type': 'h4', 'text': stripped[5:], 'level': 4, 'items': []})
        # Lists
        elif stripped.startswith('- ') or stripped.startswith('* '):
            items = [stripped[2:]]
            j = i + 1
            while j < len(lines):
                lj = lines[j].strip()
                if lj.startswith('- ') or lj.startswith('* '):
                    items.append(lj[2:])
                    j += 1
                elif not lj:
                    j += 1
                else:
                    break
            tokens.append({'type': 'ul', 'text': '', 'level': 0, 'items': items})
            i = j
            continue
        elif re.match(r'^\d+\.\s', stripped):
            items = [re.sub(r'^\d+\.\s*', '', stripped)]
            j = i + 1
            while j < len(lines):
                lj = lines[j].strip()
                if re.match(r'^\d+\.\s', lj):
                    items.append(re.sub(r'^\d+\.\s*', '', lj))
                    j += 1
                elif not lj:
                    j += 1
                else:
                    break
            tokens.append({'type': 'ol', 'text': '', 'level': 0, 'items': items})
            i = j
            continue
        # Image: ![alt](url)
        img_match = re.match(r'^!\[([^\]]*)\]\(([^)]+)\)$', stripped)
        if img_match:
            alt, src = img_match.group(1), img_match.group(2)
            tokens.append({'type': 'image', 'text': alt or 'image', 'level': 0, 'items': [], 'src': src})
            i += 1
            continue
        # Paragraph or blank
        if stripped:
            tokens.append({'type': 'p', 'text': stripped, 'level': 0, 'items': []})
        else:
            tokens.append({'type': 'blank', 'text': '', 'level': 0, 'items': []})
        i += 1
    return tokens


# ---- Image Upload ----

@router.post("/documents/{doc_id}/images")
def upload_image(
    doc_id: str,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    """Upload an image for use in markdown content."""
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp"):
        raise HTTPException(status_code=400, detail=f"Unsupported image type: {ext}")

    content = file.file.read()
    if len(content) > 10 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="Image too large (max 10MB)")

    save_dir = _ensure_dir(os.path.join(UPLOAD_DIR, "documents", doc_id, "images"))
    stored_name = str(uuid_lib.uuid4()) + ext
    save_path = os.path.join(save_dir, stored_name)

    with open(save_path, "wb") as f:
        f.write(content)

    # Return the image URL
    return {"ok": True, "url": f"/api/sops/documents/{doc_id}/images/{stored_name}"}


@router.get("/documents/{doc_id}/images/{filename}")
def get_image(doc_id: str, filename: str):
    """Serve an uploaded image."""
    file_path = os.path.join(UPLOAD_DIR, "documents", doc_id, "images", filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Image not found")
    ext = os.path.splitext(filename)[1].lower()
    mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp", ".svg": "image/svg+xml", ".bmp": "image/bmp"}
    return FileResponse(file_path, media_type=mime_map.get(ext, "image/png"))


# ---- Documents CRUD ----

@router.get("/documents", response_model=list[SopDocumentOut])
def list_documents(
    category: str | None = None,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "view")),
    current_user: User = Depends(get_current_user),
):
    query = db.query(SopDocument).filter(SopDocument.is_deleted == False)
    if current_user.role != "admin":
        query = query.filter(SopDocument.created_by == current_user.id)
    if category:
        query = query.filter(SopDocument.category == category)
    return query.order_by(SopDocument.updated_at.desc()).all()


@router.post("/documents", response_model=SopDocumentOut, status_code=201)
def create_document(
    data: SopDocumentCreate,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    doc = SopDocument(**data.model_dump(), created_by=current_user.id)
    db.add(doc)
    db.commit()
    db.refresh(doc)
    log_operation(db, current_user, "create", "sop_document", doc.name)
    return SopDocumentOut.model_validate(doc)


@router.put("/documents/{doc_id}", response_model=SopDocumentOut)
def update_document(
    doc_id: str,
    data: SopDocumentUpdate,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "edit")),
    current_user: User = Depends(get_current_user),
):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    update_data = data.model_dump(exclude_unset=True)
    details_parts = []
    for key, value in update_data.items():
        old_val = getattr(doc, key, None)
        if key == "content":
            details_parts.append("内容: 已编辑" if old_val else "内容: 已填写")
        elif key == "file_json":
            details_parts.append("附件: 已更新")
        elif key == "name":
            details_parts.append(f"名称: {old_val} → {value}")
        elif key == "file_json" and isinstance(value, dict):
            value = json.dumps(value)
        setattr(doc, key, value)
    db.commit()
    db.refresh(doc)
    log_operation(db, current_user, "update", "sop_document", doc.name,
                  details="; ".join(details_parts) if details_parts else None)
    return SopDocumentOut.model_validate(doc)


@router.delete("/documents/{doc_id}")
def delete_document(
    doc_id: str,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "delete")),
    current_user: User = Depends(get_current_user),
):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    doc.is_deleted = True
    doc.deleted_at = datetime.utcnow()
    doc.deleted_by = current_user.display_name or current_user.username
    db.commit()
    log_operation(db, current_user, "delete", "sop_document", doc.name)
    return {"ok": True}


@router.post("/documents/{doc_id}/upload")
def upload_document_file(
    doc_id: str,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    allowed = ALLOWED_REPORT_EXT if doc.category == "report" else ALLOWED_DOC_EXT
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in allowed:
        raise HTTPException(status_code=400, detail=f"File type {ext} not allowed")

    save_dir = _ensure_dir(os.path.join(UPLOAD_DIR, "documents", doc_id))
    stored_name = str(uuid_lib.uuid4()) + ext
    save_path = os.path.join(save_dir, stored_name)

    content = file.file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 50MB)")

    with open(save_path, "wb") as f:
        f.write(content)

    metadata = json.dumps({
        "original_filename": file.filename,
        "stored_filename": stored_name,
        "size": len(content),
        "uploaded_at": datetime.utcnow().isoformat(),
    })

    doc.file_json = metadata
    db.commit()
    return {"ok": True, "file": json.loads(metadata)}


@router.get("/documents/{doc_id}/download")
def download_document_file(
    doc_id: str,
    inline: bool = Query(False),
    db: Session = Depends(get_db),
):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc or not doc.file_json:
        raise HTTPException(status_code=404, detail="No file uploaded")

    try:
        meta = json.loads(doc.file_json) if isinstance(doc.file_json, str) else doc.file_json
    except json.JSONDecodeError:
        raise HTTPException(status_code=404, detail="File metadata corrupted")

    file_path = os.path.join(UPLOAD_DIR, "documents", doc_id, meta["stored_filename"])
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on disk")

    from urllib.parse import quote
    filename = meta.get("original_filename", "download")
    ext = os.path.splitext(filename)[1].lower()
    MIME = {
        ".pdf": "application/pdf", ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    disposition = "inline" if inline else "attachment"
    return FileResponse(
        file_path,
        media_type=MIME.get(ext, "application/octet-stream"),
        headers={"Content-Disposition": f"{disposition}; filename*=UTF-8''{quote(filename, safe='')}"},
    )


# ---- Document Export (MD / DOCX / PDF with rendered formatting) ----

@router.get("/documents/{doc_id}/export")
def export_document(
    doc_id: str,
    format: str = Query("md"),
    db: Session = Depends(get_db),
):
    # No auth required — doc IDs are UUIDs (unguessable), matching file download pattern
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    content = doc.content or ""
    name = doc.name or "document"
    from urllib.parse import quote

    if format == "md":
        from fastapi.responses import PlainTextResponse
        return PlainTextResponse(
            content,
            media_type="text/markdown",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}.md"},
        )

    tokens = _render_markdown(content)

    if format == "docx":
        try:
            from docx import Document as DocxDoc
            from docx.shared import Pt, Inches
            from docx.oxml import OxmlElement
            from io import BytesIO
            from fastapi.responses import StreamingResponse
        except ImportError:
            raise HTTPException(status_code=500, detail="python-docx not installed")

        word = DocxDoc()
        style = word.styles['Normal']
        style.font.size = Pt(11)
        style.font.name = '微软雅黑'
        rpr = style.element.get_or_add_rPr()
        rFonts = rpr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts')
        if rFonts is None:
            rFonts = OxmlElement('w:rFonts')
            rpr.insert(0, rFonts)
        rFonts.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia', '微软雅黑')

        def _font(run, name='微软雅黑', size=11, bold=False, italic=False):
            run.font.name = name
            run.font.size = Pt(size)
            run.bold = bold
            run.italic = italic
            rpr = run._element.get_or_add_rPr()
            rFonts = rpr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts')
            if rFonts is None:
                rFonts = OxmlElement('w:rFonts')
                rpr.insert(0, rFonts)
            rFonts.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia', '微软雅黑')

        def _add(text, size=11, bold=False, italic=False):
            p = word.add_paragraph()
            run = p.add_run(text)
            _font(run, size=size, bold=bold, italic=italic)
            return p

        def _set_cell_font(cell, size=9):
            for para in cell.paragraphs:
                for run in para.runs:
                    _font(run, size=size)
                if not para.runs:
                    # Text set via cell.text — re-add as run
                    text = para.text
                    para.clear()
                    run = para.add_run(text)
                    _font(run, size=size)

        for t in tokens:
            tp = t['type']
            if tp == 'blank':
                word.add_paragraph("")
            elif tp in ('h1',):
                _add(t['text'], size=18, bold=True)
            elif tp in ('h2',):
                _add(t['text'], size=15, bold=True)
            elif tp in ('h3',):
                _add(t['text'], size=13, bold=True)
            elif tp in ('h4', 'h5', 'h6'):
                _add(t['text'], size=12, bold=True)
            elif tp == 'ul':
                for item in t['items']:
                    _add(item, size=11)
            elif tp == 'ol':
                for item in t['items']:
                    _add(item, size=11)
            elif tp == 'code':
                p = word.add_paragraph()
                run = p.add_run(t['text'])
                _font(run, name='Courier New', size=9)
            elif tp == 'blockquote':
                p = word.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.5)
                run = p.add_run(t['text'])
                _font(run, size=11, italic=True)
            elif tp == 'table':
                rows = t.get('items', [])
                if rows:
                    table = word.add_table(rows=len(rows), cols=len(rows[0]) if rows else 1)
                    table.style = 'Light Grid Accent 1'
                    for ri, row in enumerate(rows):
                        for ci, cell_text in enumerate(row):
                            if ci < len(table.rows[ri].cells):
                                cell = table.rows[ri].cells[ci]
                                cell.text = ''
                                run = cell.paragraphs[0].add_run(cell_text)
                                _font(run, size=9)
                    word.add_paragraph("")
            elif tp == 'image':
                src = t.get('src', '')
                img_path = _resolve_image_path(src)
                if img_path:
                    try:
                        p = word.add_paragraph()
                        run = p.add_run()
                        run.add_picture(img_path, width=Inches(5.5))
                        p.alignment = 1  # center
                    except Exception:
                        _add(f'[图片: {t["text"]}]', size=10)
                else:
                    _add(f'[图片: {t["text"]}]', size=10)
            elif tp == 'p':
                _add(t['text'], size=11)

        out = BytesIO()
        word.save(out); out.seek(0)
        return StreamingResponse(
            out,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}.docx"},
        )

    if format == "pdf":
        try:
            from fpdf import FPDF
            from io import BytesIO
            from fastapi.responses import StreamingResponse
        except ImportError:
            raise HTTPException(status_code=500, detail="fpdf2 not installed")

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(True, 20)

        if CJK_FONT_PATH:
            pdf.add_font("CJK", "", CJK_FONT_PATH, uni=True)
            fn = "CJK"
        else:
            fn = "Helvetica"

        for t in tokens:
            tp = t['type']
            pdf.set_font(fn, "", 11)
            if tp == 'blank':
                pdf.ln(4)
            elif tp == 'h1':
                pdf.set_font(fn, "", 18)
                pdf.multi_cell(0, 10, t['text'], new_x="LMARGIN", new_y="NEXT")
                pdf.ln(4)
            elif tp == 'h2':
                pdf.set_font(fn, "", 15)
                pdf.multi_cell(0, 8, t['text'], new_x="LMARGIN", new_y="NEXT")
                pdf.ln(2)
            elif tp == 'h3':
                pdf.set_font(fn, "", 13)
                pdf.multi_cell(0, 7, t['text'], new_x="LMARGIN", new_y="NEXT")
                pdf.ln(2)
            elif tp in ('h4', 'h5', 'h6'):
                pdf.set_font(fn, "", 12)
                pdf.multi_cell(0, 7, t['text'], new_x="LMARGIN", new_y="NEXT")
            elif tp == 'ul':
                pdf.set_font(fn, "", 11)
                for item in t['items']:
                    pdf.set_x(14)
                    pdf.multi_cell(186, 5.5, f"• {item}", new_x="LMARGIN", new_y="NEXT")
            elif tp == 'ol':
                pdf.set_font(fn, "", 11)
                for idx, item in enumerate(t['items'], 1):
                    pdf.set_x(14)
                    pdf.multi_cell(186, 5.5, f"{idx}. {item}", new_x="LMARGIN", new_y="NEXT")
            elif tp == 'code':
                pdf.set_font(fn, "", 9)
                for cl in t['text'].split('\n'):
                    pdf.set_x(14)
                    pdf.multi_cell(186, 4.5, cl, new_x="LMARGIN", new_y="NEXT")
                pdf.ln(2)
            elif tp == 'blockquote':
                pdf.set_font(fn, "", 11)
                pdf.set_x(16)
                pdf.set_font(fn, "", 10)
                pdf.multi_cell(180, 5.5, t['text'], new_x="LMARGIN", new_y="NEXT")
            elif tp == 'table':
                rows = t.get('items', [])
                if rows:
                    # Normalize all rows to same column count
                    cols = max(len(r) for r in rows)
                    norm_rows = []
                    for r in rows:
                        nr = list(r)
                        while len(nr) < cols:
                            nr.append('')
                        norm_rows.append(nr)
                    # Fixed-width columns for narrow content, proportional for wider tables
                    page_w = 190
                    if cols <= 4:
                        col_w = page_w / cols
                    else:
                        col_w = max(22, page_w / cols)
                    # Calculate row heights
                    row_heights = []
                    for row in norm_rows:
                        max_h = 6
                        for cell_text in row:
                            text_w = pdf.get_string_width(cell_text)
                            lines = max(1, int(text_w / (col_w - 1.5)) + 1)
                            max_h = max(max_h, lines * 4.5 + 1)
                        row_heights.append(max_h)
                    # Check if table fits on page
                    total_h = sum(row_heights)
                    if pdf.get_y() + total_h > 270:
                        pdf.add_page()
                    # Draw the table
                    for ri, row in enumerate(norm_rows):
                        rh = row_heights[ri]
                        y0 = pdf.get_y()
                        # Check page break for this row
                        if y0 + rh > 270:
                            pdf.add_page()
                            y0 = pdf.get_y()
                        # Header background
                        if ri == 0:
                            pdf.set_fill_color(230, 235, 245)
                            pdf.set_font(fn, "", 8)
                        else:
                            pdf.set_fill_color(255, 255, 255)
                            pdf.set_font(fn, "", 8)
                        # Draw cell backgrounds and borders, then text
                        for ci, cell_text in enumerate(row):
                            x0 = 10 + ci * col_w
                            # Cell rect
                            pdf.set_fill_color(230 if ri == 0 else 255, 235 if ri == 0 else 255, 245 if ri == 0 else 255)
                            pdf.rect(x0, y0, col_w, rh, style='DF')
                            # Cell text — clipped within cell bounds
                            pdf.set_xy(x0 + 0.5, y0 + 0.5)
                            # Manually break text into lines that fit
                            words = cell_text
                            line_h = 4.5
                            max_w = col_w - 1.5
                            # Simple word-wrap by character
                            current_line = ''
                            for ch in words:
                                test_line = current_line + ch
                                if pdf.get_string_width(test_line) > max_w and current_line:
                                    pdf.set_x(x0 + 0.5)
                                    pdf.cell(max_w, line_h, current_line, new_x="LMARGIN", new_y="NEXT")
                                    current_line = ch
                                else:
                                    current_line = test_line
                            if current_line:
                                pdf.set_x(x0 + 0.5)
                                pdf.cell(max_w, line_h, current_line, new_x="LMARGIN", new_y="NEXT")
                        # Move to next row position
                        pdf.set_xy(10, y0 + rh)
                    pdf.ln(4)
            elif tp == 'image':
                src = t.get('src', '')
                img_path = _resolve_image_path(src)
                if img_path:
                    try:
                        # Fit image within page width
                        from PIL import Image as PILImage
                        with PILImage.open(img_path) as pil_img:
                            pw, ph = pil_img.size
                        max_w = 190  # mm
                        max_h = 120  # mm
                        ratio = min(max_w / pw, max_h / ph, 1.0)
                        w, h = pw * ratio, ph * ratio
                        if pdf.get_y() + h > 260:
                            pdf.add_page()
                        pdf.image(img_path, x=10, y=pdf.get_y(), w=w, h=h)
                        pdf.set_y(pdf.get_y() + h + 4)
                    except ImportError:
                        pdf.set_font(fn, "", 10)
                        pdf.cell(0, 6, f'[图片: {t["text"]}]', new_x="LMARGIN", new_y="NEXT")
                    except Exception:
                        pdf.set_font(fn, "", 10)
                        pdf.cell(0, 6, f'[图片: {t["text"]}]', new_x="LMARGIN", new_y="NEXT")
                else:
                    pdf.set_font(fn, "", 10)
                    pdf.cell(0, 6, f'[图片: {t["text"]}]', new_x="LMARGIN", new_y="NEXT")
            elif tp == 'p':
                pdf.set_font(fn, "", 11)
                pdf.multi_cell(0, 6, t['text'], new_x="LMARGIN", new_y="NEXT")

        out = BytesIO()
        pdf.output(out); out.seek(0)
        return StreamingResponse(
            out,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(name)}.pdf"},
        )

    raise HTTPException(status_code=400, detail=f"Unknown format: {format}")


# ---- PPT/PDF Preview ----

@router.get("/documents/{doc_id}/preview")
def preview_document(
    doc_id: str,
    db: Session = Depends(get_db),
):
    """Convert uploaded PPT to PDF for preview using python-pptx + fpdf2."""
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id, SopDocument.is_deleted == False).first()
    if not doc or not doc.file_json:
        raise HTTPException(status_code=404, detail="No file uploaded")

    meta = json.loads(doc.file_json) if isinstance(doc.file_json, str) else doc.file_json
    file_path = os.path.join(UPLOAD_DIR, "documents", doc_id, meta["stored_filename"])
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on disk")

    filename = meta.get("original_filename", "download")
    ext = os.path.splitext(filename)[1].lower()

    from urllib.parse import quote

    # PDF: serve directly
    if ext == ".pdf":
        return FileResponse(
            file_path,
            media_type="application/pdf",
            headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(filename, safe='')}"},
        )

    # PPT/PPTX: convert to PDF
    if ext in (".ppt", ".pptx"):
        from pptx import Presentation
        from fpdf import FPDF
        from io import BytesIO
        from fastapi.responses import StreamingResponse

        prs = Presentation(file_path)

        pdf = FPDF(orientation='L', unit='mm', format='A4')
        if CJK_FONT_PATH:
            pdf.add_font("CJK", "", CJK_FONT_PATH, uni=True)
            fn = "CJK"
        else:
            fn = "Helvetica"
        pdf.set_auto_page_break(True, 10)

        for si, slide in enumerate(prs.slides, 1):
            pdf.add_page()
            pdf.set_font(fn, "", 16)
            pdf.cell(0, 10, f"Slide {si}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(6)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    pdf.set_font(fn, "", 10)
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            pdf.multi_cell(0, 5.5, text, new_x="LMARGIN", new_y="NEXT")
                    pdf.ln(2)

        out = BytesIO()
        pdf.output(out); out.seek(0)
        return StreamingResponse(
            out,
            media_type="application/pdf",
            headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(os.path.splitext(filename)[0])}.pdf"},
        )

    raise HTTPException(status_code=400, detail="Preview not available for this file type")


# ---- Test Case Categories ----

@router.get("/test-case-categories", response_model=list[dict])
def list_test_case_categories(
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "view")),
    current_user: User = Depends(get_current_user),
):
    query = db.query(TestCaseCategory).filter(TestCaseCategory.is_deleted == False)
    if current_user.role != "admin":
        query = query.filter(TestCaseCategory.created_by == current_user.id)
    cats = query.order_by(TestCaseCategory.name).all()
    # Count test cases per category
    result = []
    for c in cats:
        count = db.query(TestCase).filter(
            TestCase.category_id == c.id, TestCase.is_deleted == False
        ).count()
        result.append({"id": c.id, "name": c.name, "case_count": count})
    return result


@router.post("/test-case-categories", status_code=201)
def create_test_case_category(
    data: dict,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    name = data.get("name", "").strip()
    if not name:
        raise HTTPException(status_code=400, detail="Name required")
    existing = db.query(TestCaseCategory).filter(
        TestCaseCategory.name == name, TestCaseCategory.is_deleted == False
    ).first()
    if existing:
        raise HTTPException(status_code=400, detail="Category already exists")
    cat = TestCaseCategory(name=name, created_by=current_user.id)
    db.add(cat)
    db.commit()
    db.refresh(cat)
    log_operation(db, current_user, "create", "test_case_category", name)
    return {"id": cat.id, "name": cat.name, "case_count": 0}


@router.put("/test-case-categories/{cat_id}")
def update_test_case_category(
    cat_id: str,
    data: dict,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "edit")),
    current_user: User = Depends(get_current_user),
):
    cat = db.query(TestCaseCategory).filter(TestCaseCategory.id == cat_id, TestCaseCategory.is_deleted == False).first()
    if not cat:
        raise HTTPException(status_code=404, detail="Category not found")
    name = data.get("name", "").strip()
    if name:
        cat.name = name
    db.commit()
    return {"ok": True}


@router.delete("/test-case-categories/{cat_id}")
def delete_test_case_category(
    cat_id: str,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "delete")),
    current_user: User = Depends(get_current_user),
):
    cat = db.query(TestCaseCategory).filter(TestCaseCategory.id == cat_id, TestCaseCategory.is_deleted == False).first()
    if not cat:
        raise HTTPException(status_code=404, detail="Category not found")
    # Unlink test cases
    db.query(TestCase).filter(TestCase.category_id == cat_id).update({TestCase.category_id: None})
    cat.is_deleted = True
    cat.deleted_at = datetime.utcnow()
    cat.deleted_by = current_user.display_name or current_user.username
    db.commit()
    log_operation(db, current_user, "delete", "test_case_category", cat.name)
    return {"ok": True}


# ---- Test Cases ----

@router.get("/test-cases", response_model=TestCaseListOut)
def list_test_cases(
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=500),
    category_id: str | None = None,
    module: str | None = None,
    priority: str | None = None,
    status: str | None = None,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "view")),
    current_user: User = Depends(get_current_user),
):
    query = db.query(TestCase).filter(TestCase.is_deleted == False)
    if current_user.role != "admin":
        query = query.filter(TestCase.created_by == current_user.id)
    if category_id:
        query = query.filter(TestCase.category_id == category_id)
    if module:
        query = query.filter(TestCase.module == module)
    if priority:
        query = query.filter(TestCase.priority == priority)
    if status:
        query = query.filter(TestCase.status == status)

    total = query.count()
    items = query.order_by(TestCase.updated_at.desc()).offset((page - 1) * page_size).limit(page_size).all()
    return TestCaseListOut(
        items=[TestCaseOut.model_validate(i) for i in items],
        total=total, page=page, page_size=page_size,
    )


@router.post("/test-cases", response_model=TestCaseOut, status_code=201)
def create_test_case(
    data: TestCaseCreate,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    tc = TestCase(**data.model_dump(), created_by=current_user.id)
    db.add(tc)
    db.commit()
    db.refresh(tc)
    log_operation(db, current_user, "create", "test_case", tc.title)
    return TestCaseOut.model_validate(tc)


@router.put("/test-cases/{tc_id}", response_model=TestCaseOut)
def update_test_case(
    tc_id: str,
    data: TestCaseUpdate,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "edit")),
    current_user: User = Depends(get_current_user),
):
    tc = db.query(TestCase).filter(TestCase.id == tc_id, TestCase.is_deleted == False).first()
    if not tc:
        raise HTTPException(status_code=404, detail="Test case not found")
    field_labels = {"title": "标题", "module": "模块", "priority": "优先级",
                    "precondition": "前置条件", "steps": "测试步骤",
                    "expected_result": "预期结果", "status": "状态", "remarks": "备注",
                    "category_id": "所属客户端"}
    changes = []
    update_data = data.model_dump(exclude_unset=True)
    for key, new_val in update_data.items():
        old_val = getattr(tc, key, None)
        if old_val != new_val:
            label = field_labels.get(key, key)
            changes.append(f"{label}: {old_val} → {new_val}")
    for key, value in update_data.items():
        setattr(tc, key, value)
    db.commit()
    db.refresh(tc)
    log_operation(db, current_user, "update", "test_case", tc.title,
                  details="; ".join(changes) if changes else None)
    return TestCaseOut.model_validate(tc)


@router.delete("/test-cases/{tc_id}")
def delete_test_case(
    tc_id: str,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "delete")),
    current_user: User = Depends(get_current_user),
):
    tc = db.query(TestCase).filter(TestCase.id == tc_id, TestCase.is_deleted == False).first()
    if not tc:
        raise HTTPException(status_code=404, detail="Test case not found")
    tc.is_deleted = True
    tc.deleted_at = datetime.utcnow()
    tc.deleted_by = current_user.display_name or current_user.username
    db.commit()
    log_operation(db, current_user, "delete", "test_case", tc.title)
    return {"ok": True}


@router.post("/test-cases/import")
def import_test_cases(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in (".xlsx", ".xls"):
        raise HTTPException(status_code=400, detail="Only .xlsx files are supported")

    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl not installed on server")

    try:
        from io import BytesIO
        contents = BytesIO(file.file.read())
        wb = openpyxl.load_workbook(contents, read_only=True)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法解析文件: {str(e)}")
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    if len(rows) < 2:
        return {"ok": True, "count": 0}

    headers = [str(h or "").strip() for h in rows[0]]
    col_map = {
        "title": "title", "module": "module", "priority": "priority",
        "precondition": "precondition", "steps": "steps", "expected_result": "expected_result",
        "status": "status", "remarks": "remarks",
        # Chinese headers (from export)
        "标题": "title", "客户端": "category", "模块": "module", "优先级": "priority",
        "前置条件": "precondition", "测试步骤": "steps", "预期结果": "expected_result",
        "状态": "status", "备注": "remarks",
    }
    idx_map = {}
    for i, h in enumerate(headers):
        key = col_map.get(h)
        if key:
            idx_map[i] = key

    # Reverse status mapping
    rev_status = {"草稿": "draft", "就绪": "ready", "废弃": "deprecated", "draft": "draft", "ready": "ready", "deprecated": "deprecated"}

    try:
        count = 0
        for row in rows[1:]:
            if all(v is None for v in row):
                continue
            values: dict[str, Any] = {}
            category_name = None
            for i, field in idx_map.items():
                val = row[i] if i < len(row) else None
                if val is not None:
                    v = str(val).strip()
                    if field == "category":
                        category_name = v
                    elif field == "status":
                        values["status"] = rev_status.get(v, "draft")
                    else:
                        values[field] = v
            if not values.get("title"):
                continue
            values.setdefault("priority", "P2")
            values.setdefault("status", "draft")

            # Resolve category name → category_id
            if category_name:
                cat = db.query(TestCaseCategory).filter(TestCaseCategory.name == category_name).first()
                if not cat:
                    cat = TestCaseCategory(name=category_name, created_by=current_user.id)
                    db.add(cat)
                    db.flush()
                values["category_id"] = cat.id

            tc = TestCase(**values, created_by=current_user.id)
            db.add(tc)
            count += 1

        db.commit()
        log_operation(db, current_user, "import", "test_case", f"{count} cases")
        return {"ok": True, "count": count}
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=400, detail=f"导入失败: {str(e)}")


@router.get("/test-cases/template")
def download_test_case_template():
    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl not installed on server")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "案例模板"
    headers = ["标题", "客户端", "模块", "优先级", "前置条件", "测试步骤", "预期结果", "状态", "备注"]
    ws.append(headers)
    # Add example row
    ws.append(["示例：登录功能测试", "客户端名称", "登录模块", "P1", "用户已注册", "1.打开登录页\n2.输入账号密码\n3.点击登录", "成功登录并跳转首页", "就绪", ""])

    from io import BytesIO
    from fastapi.responses import StreamingResponse
    from urllib.parse import quote
    output = BytesIO()
    wb.save(output); output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote('案例导入模板.xlsx', safe='')}"},
    )


@router.get("/test-cases/export")
def export_test_cases(
    category_id: str | None = None,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "view")),
    current_user: User = Depends(get_current_user),
):
    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl not installed on server")

    query = db.query(TestCase).filter(TestCase.is_deleted == False)
    if current_user.role != "admin":
        query = query.filter(TestCase.created_by == current_user.id)
    if category_id:
        query = query.filter(TestCase.category_id == category_id)
    items = query.order_by(TestCase.updated_at.desc()).all()

    # Resolve category names
    cat_names: dict[str, str] = {}
    cat_ids = {tc.category_id for tc in items if tc.category_id}
    if cat_ids:
        cats = db.query(TestCaseCategory).filter(TestCaseCategory.id.in_(cat_ids)).all()
        cat_names = {c.id: c.name for c in cats}

    # Status mapping to Chinese
    status_map = {"draft": "草稿", "ready": "就绪", "deprecated": "废弃"}

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "test_cases"
    # Same order as frontend display: 标题, 客户端, 模块, 优先级, 前置条件, 测试步骤, 预期结果, 状态, 备注
    headers = ["标题", "客户端", "模块", "优先级", "前置条件", "测试步骤", "预期结果", "状态", "备注"]
    ws.append(headers)
    for tc in items:
        ws.append([tc.title, cat_names.get(tc.category_id or '', ''), tc.module, tc.priority, tc.precondition, tc.steps, tc.expected_result, status_map.get(tc.status, tc.status), tc.remarks])

    from io import BytesIO
    from fastapi.responses import StreamingResponse
    from urllib.parse import quote
    output = BytesIO()
    wb.save(output); output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote('test_cases_export.xlsx', safe='')}"},
    )


# ---- Scripts ----

@router.get("/scripts", response_model=list[ScriptOut])
def list_scripts(
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "view")),
    current_user: User = Depends(get_current_user),
):
    query = db.query(ScriptFile).filter(ScriptFile.is_deleted == False)
    if current_user.role != "admin":
        query = query.filter(ScriptFile.created_by == current_user.id)
    return query.order_by(ScriptFile.updated_at.desc()).all()


@router.post("/scripts", response_model=ScriptOut, status_code=201)
def create_script(
    name: str = Query(...),
    description: str | None = Query(None),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "create")),
    current_user: User = Depends(get_current_user),
):
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext != ".zip":
        raise HTTPException(status_code=400, detail="Only .zip files are supported")

    script = ScriptFile(name=name, description=description, created_by=current_user.id)
    db.add(script)
    db.flush()

    save_dir = _ensure_dir(os.path.join(UPLOAD_DIR, "scripts", script.id))
    stored_name = str(uuid_lib.uuid4()) + ext
    save_path = os.path.join(save_dir, stored_name)

    content = file.file.read()
    if len(content) > MAX_FILE_SIZE:
        db.rollback()
        raise HTTPException(status_code=400, detail="File too large (max 50MB)")

    with open(save_path, "wb") as f:
        f.write(content)

    script.file_json = json.dumps({
        "original_filename": file.filename,
        "stored_filename": stored_name,
        "size": len(content),
        "uploaded_at": datetime.utcnow().isoformat(),
    })
    db.commit()
    db.refresh(script)
    log_operation(db, current_user, "create", "script", script.name)
    return ScriptOut.model_validate(script)


@router.delete("/scripts/{script_id}")
def delete_script(
    script_id: str,
    db: Session = Depends(get_db),
    _=Depends(require_permission("sop", "delete")),
    current_user: User = Depends(get_current_user),
):
    script = db.query(ScriptFile).filter(ScriptFile.id == script_id, ScriptFile.is_deleted == False).first()
    if not script:
        raise HTTPException(status_code=404, detail="Script not found")
    script.is_deleted = True
    script.deleted_at = datetime.utcnow()
    script.deleted_by = current_user.display_name or current_user.username
    db.commit()
    log_operation(db, current_user, "delete", "script", script.name)
    return {"ok": True}


@router.get("/scripts/{script_id}/download")
def download_script(
    script_id: str,
    db: Session = Depends(get_db),
):
    script = db.query(ScriptFile).filter(ScriptFile.id == script_id, ScriptFile.is_deleted == False).first()
    if not script or not script.file_json:
        raise HTTPException(status_code=404, detail="No file uploaded")

    meta = json.loads(script.file_json) if isinstance(script.file_json, str) else script.file_json
    file_path = os.path.join(UPLOAD_DIR, "scripts", script_id, meta["stored_filename"])
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on disk")

    from urllib.parse import quote
    filename = meta.get("original_filename", "script.zip")
    return FileResponse(
        file_path,
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename, safe='')}"},
    )
